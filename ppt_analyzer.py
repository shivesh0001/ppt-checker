#!/usr/bin/env python3

import json
import time
import re
from typing import List
from models import Inconsistency, SlideContent

try:
    from pptx import Presentation
    import google.generativeai as genai
    import pytesseract
except ImportError as e:
    print(f"Error: Missing required dependency - {e}")
    print("Please install dependencies: pip install python-pptx google-generativeai pillow pytesseract")
    exit(1)

class PPTAnalyzer:
    
    def __init__(self, api_key: str, batch_size: int = 6):
        self.api_key = api_key
        self.batch_size = batch_size
        self.model = None
        self._setup_gemini()
    
    def _setup_gemini(self):
        # Setup the Gemini AI client
        try:
            genai.configure(api_key=self.api_key)
            self.model = genai.GenerativeModel('gemini-2.0-flash-exp')
        except Exception as e:
            raise Exception(f"Failed to setup Gemini API: {e}")
    
    def extract_slide_content(self, pptx_path: str, use_ocr: bool = False) -> List[SlideContent]:
        # Extract all text from the PowerPoint file
        try:
            presentation = Presentation(pptx_path)
            slides = []
            
            for i, slide in enumerate(presentation.slides, 1):
                text = self._get_slide_text(slide)
                ocr_text = ""
                
                if use_ocr:
                    try:
                        ocr_text = self._get_ocr_text(slide)
                    except:
                        print(f"Warning: OCR failed for slide {i}")
                
                slides.append(SlideContent(
                    slide_number=i,
                    text=text,
                    ocr_text=ocr_text
                ))
                
            return slides
            
        except Exception as e:
            raise Exception(f"Failed to extract content from PowerPoint: {e}")
    
    def _get_slide_text(self, slide) -> str:
        # Pull out all the text from a slide
        text_parts = []
        
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text_parts.append(shape.text.strip())
            elif hasattr(shape, "text_frame") and shape.text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    para_text = ""
                    for run in paragraph.runs:
                        para_text += run.text
                    if para_text.strip():
                        text_parts.append(para_text.strip())
        
        return "\n".join(text_parts)
    
    def _get_ocr_text(self, slide) -> str:
        # Try to extract text from images using OCR
        ocr_text = []
        
        for shape in slide.shapes:
            if hasattr(shape, 'image'):
                try:
                    # TODO: implement proper image OCR extraction
                    # For now just skip it
                    pass
                except:
                    continue
        
        return "\n".join(ocr_text)
    
    def analyze_inconsistencies(self, slides: List[SlideContent]) -> List[Inconsistency]:
        print(f"Analyzing {len(slides)} slides for inconsistencies...")
        
        # First pass: check slides in small batches
        batch_results = self._batch_analysis(slides)
        
        # Second pass: look for issues across all slides  
        global_results = self._cross_slide_analysis(slides)
        
        # Combine everything and filter out low confidence findings
        all_issues = batch_results + global_results
        good_issues = [issue for issue in all_issues if issue.confidence >= 0.7]
        
        # Clean up duplicates
        final_issues = self._remove_duplicates(good_issues)
        
        return final_issues
    
    def _batch_analysis(self, slides: List[SlideContent]) -> List[Inconsistency]:
        issues = []
        
        for i in range(0, len(slides), self.batch_size):
            batch = slides[i:i + self.batch_size]
            batch_num = i//self.batch_size + 1
            print(f"Processing batch {batch_num} ({len(batch)} slides)...")
            
            try:
                batch_issues = self._analyze_batch(batch)
                issues.extend(batch_issues)
                time.sleep(1)  # don't hammer the API
            except:
                print(f"Warning: Batch {batch_num} analysis failed")
                continue
        
        return issues
    
    def _cross_slide_analysis(self, slides: List[SlideContent]) -> List[Inconsistency]:
        print("Checking for issues across all slides...")
        
        # Create summary of important numbers/dates from all slides
        summary = self._create_slide_summary(slides)
        
        try:
            return self._analyze_global_stuff(summary, slides)
        except:
            print("Warning: Global analysis failed")
            return []
    
    def _analyze_batch(self, batch: List[SlideContent]) -> List[Inconsistency]:        
        # Format slides for the AI to analyze
        batch_text = self._format_slides(batch)
        prompt = self._build_prompt(batch_text, is_batch=True)
        
        try:
            response = self.model.generate_content(prompt)
            return self._parse_response(response.text, batch)
        except:
            print("API call failed")
            return []
    
    def _analyze_global_stuff(self, summary: str, slides: List[SlideContent]) -> List[Inconsistency]:
        prompt = self._build_prompt(summary, is_batch=False)
        
        try:
            response = self.model.generate_content(prompt)
            return self._parse_response(response.text, slides)
        except:
            print("Global analysis API call failed")
            return []
    
    def _format_slides(self, batch: List[SlideContent]) -> str:
        parts = []
        
        for slide in batch:
            content = f"=== SLIDE {slide.slide_number} ===\n"
            content += slide.text
            if slide.ocr_text:
                content += f"\n[OCR TEXT]: {slide.ocr_text}"
            parts.append(content)
        
        return "\n\n".join(parts)
    
    def _create_slide_summary(self, slides: List[SlideContent]) -> str:
        # Extract important numbers and dates from each slide
        summary_parts = []
        
        for slide in slides:
            # Look for numbers like $100M, 25%, etc
            numbers = re.findall(r'\$?[\d,]+\.?\d*[%MBK]?', slide.text)
            # Look for dates in various formats
            dates = re.findall(r'\b\d{4}\b|\b\d{1,2}\/\d{1,2}\/\d{2,4}\b|\b\w+\s+\d{1,2},?\s+\d{4}\b', slide.text)
            
            if numbers or dates:
                summary_parts.append(f"Slide {slide.slide_number}: Numbers: {numbers}, Dates: {dates}")
        
        return "\n".join(summary_parts)
    
    def _build_prompt(self, content: str, is_batch: bool = True) -> str:
        scope = "within these slides" if is_batch else "across the entire presentation"
        
        prompt = f"""You are an expert business analyst reviewing a PowerPoint presentation for inconsistencies. 

CRITICAL INSTRUCTIONS:
- Only flag GENUINE business logic inconsistencies that would concern executives
- Ignore stylistic differences, synonyms, or different phrasings of the same concept
- Focus on factual/numerical conflicts, timeline errors, and logical contradictions
- Provide specific evidence with exact quotes
- Rate confidence 0.0-1.0 for each finding

INCONSISTENCY TYPES TO DETECT {scope}:
1. Numerical conflicts (revenue figures, percentages that don't add up, contradictory metrics)
2. Timeline mismatches (date conflicts, impossible sequences, chronological errors)
3. Logical contradictions (mutually exclusive statements about market, competition, etc.)
4. Data relationship errors (parts don't sum to whole, inconsistent breakdowns)

CONTENT TO ANALYZE:
{content}

Respond in JSON format only:
{{
    "inconsistencies": [
        {{
            "type": "Numerical Conflict|Timeline Mismatch|Logical Contradiction|Data Relationship Error",
            "confidence": 0.85,
            "slides": [3, 8],
            "issue": "Brief description of the inconsistency",
            "evidence": [
                "Exact quote from slide 3",
                "Exact quote from slide 8"
            ]
        }}
    ]
}}

If no genuine inconsistencies are found, return: {{"inconsistencies": []}}
"""
        return prompt
    
    def _parse_response(self, response_text: str, slides: List[SlideContent]) -> List[Inconsistency]:
        try:
            # Clean up the response text  
            text = response_text.strip()
            if text.startswith("```json"):
                text = text[7:-3]
            elif text.startswith("```"):
                text = text[3:-3]
            
            data = json.loads(text)
            results = []
            
            for item in data.get("inconsistencies", []):
                inc = Inconsistency(
                    type=item.get("type", "Unknown"),
                    confidence=float(item.get("confidence", 0.0)),
                    slides=item.get("slides", []),
                    issue=item.get("issue", ""),
                    evidence=item.get("evidence", [])
                )
                results.append(inc)
            
            return results
            
        except json.JSONDecodeError:
            print(f"Warning: Couldn't parse AI response as JSON")
            print(f"Response was: {response_text[:200]}...")
            return []
        except:
            print("Warning: Error parsing AI response")
            return []
    
    def _remove_duplicates(self, issues: List[Inconsistency]) -> List[Inconsistency]:
        unique_issues = []
        seen = set()
        
        for issue in issues:
            # Create identifier from slides and issue text
            key = (tuple(sorted(issue.slides)), issue.issue.lower()[:50])
            
            if key not in seen:
                seen.add(key)
                unique_issues.append(issue)
        
        return unique_issues
